from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ──────────────────────────────────────────────────
BG       = RGBColor(0x03, 0x07, 0x12)   # near-black
GREEN    = RGBColor(0x00, 0xFF, 0x88)
BLUE     = RGBColor(0x00, 0xD4, 0xFF)
PURPLE   = RGBColor(0xA8, 0x55, 0xF7)
ORANGE   = RGBColor(0xF9, 0x73, 0x16)
PINK     = RGBColor(0xEC, 0x48, 0x99)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
MUTED    = RGBColor(0x94, 0xA3, 0xB8)
CARD_BG  = RGBColor(0x0F, 0x17, 0x2A)
CARD2_BG = RGBColor(0x0A, 0x0F, 0x1E)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank

# ── Helpers ──────────────────────────────────────────────────
def add_slide():
    sl = prs.slides.add_slide(blank_layout)
    fill_bg(sl)
    return sl

def fill_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG

def box(slide, x, y, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def rounded_box(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(5, x, y, w, h)  # rounded rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape

def txt(slide, text, x, y, w, h, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return txb

def accent_line(slide, x, y, w, color=GREEN, h=Inches(0.04)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def badge(slide, text, x, y, color=GREEN):
    bg_shape = slide.shapes.add_shape(5, x, y, Inches(2.6), Inches(0.32))
    bg_shape.fill.solid()
    r, g, b = color[0], color[1], color[2]
    bg_shape.fill.fore_color.rgb = RGBColor(
        min(255, r + 30), min(255, g + 30), min(255, b + 30)
    )
    bg_shape.line.color.rgb = color
    bg_shape.line.width = Pt(0.8)
    bg_shape.adjustments[0] = 0.5
    tf = bg_shape.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(8)
    run.font.bold = True
    run.font.color.rgb = color
    run.font.name = "Courier New"
    return bg_shape

def section_tag(slide, text, color=GREEN):
    b = badge(slide, text, Inches(0.45), Inches(0.32), color)
    return b

def slide_title(slide, title, subtitle=None, color=GREEN):
    accent_line(slide, Inches(0.45), Inches(0.75), Inches(2.5), color)
    txt(slide, title, Inches(0.45), Inches(0.85), Inches(12.4), Inches(0.85),
        size=34, bold=True, color=WHITE, font="Calibri")
    if subtitle:
        txt(slide, subtitle, Inches(0.45), Inches(1.62), Inches(12.0), Inches(0.4),
            size=13, color=MUTED, font="Calibri")

def bullet(slide, items, x, y, w, color=MUTED, size=13, marker_color=GREEN):
    cur_y = y
    for item in items:
        # dot marker
        dot = slide.shapes.add_shape(9, x, cur_y + Inches(0.07), Inches(0.1), Inches(0.1))
        dot.fill.solid(); dot.fill.fore_color.rgb = marker_color
        dot.line.fill.background()
        txt(slide, item, x + Inches(0.18), cur_y, w - Inches(0.18), Inches(0.35),
            size=size, color=color)
        cur_y += Inches(0.36)
    return cur_y

def code_block(slide, code, x, y, w, h):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x05, 0x0D, 0x1E)
    shape.line.color.rgb = RGBColor(0x1E, 0x3A, 0x5F); shape.line.width = Pt(0.75)
    tf = shape.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run(); run.text = code
    run.font.size = Pt(10); run.font.color.rgb = GREEN
    run.font.name = "Courier New"

def formula_box(slide, formula, x, y, w, h, bg=CARD_BG, fg=BLUE):
    shape = slide.shapes.add_shape(5, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = fg; shape.line.width = Pt(1)
    shape.adjustments[0] = 0.05
    tf = shape.text_frame; tf.word_wrap = True
    tf.margin_left  = Inches(0.15); tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = formula
    run.font.size = Pt(13); run.font.color.rgb = fg
    run.font.name = "Courier New"; run.font.bold = True

def info_card(slide, title, body_lines, x, y, w, h, accent=GREEN):
    shape = slide.shapes.add_shape(5, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = accent; shape.line.width = Pt(0.8)
    shape.adjustments[0] = 0.04
    # accent top bar
    bar = slide.shapes.add_shape(1, x, y, w, Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    txt(slide, title, x + Inches(0.15), y + Inches(0.1), w - Inches(0.3), Inches(0.32),
        size=11, bold=True, color=accent, font="Calibri")
    cur = y + Inches(0.44)
    for line in body_lines:
        txt(slide, line, x + Inches(0.15), cur, w - Inches(0.3), Inches(0.32),
            size=10, color=MUTED, font="Calibri")
        cur += Inches(0.29)

# ════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ════════════════════════════════════════════════════════════
sl = add_slide()

# decorative gradient boxes
for i, (col, opa) in enumerate([(GREEN, 0.06), (BLUE, 0.04), (PURPLE, 0.03)]):
    b = sl.shapes.add_shape(9, Inches(9+i*0.6), Inches(0.5-i*0.3), Inches(3.5), Inches(3.5))
    b.fill.solid(); b.fill.fore_color.rgb = col
    b.line.fill.background()

accent_line(sl, Inches(1.0), Inches(2.5), Inches(4.5), GREEN, h=Inches(0.05))
txt(sl, "MONTE CARLO ENGINE", Inches(1.0), Inches(2.65), Inches(11), Inches(1.2),
    size=52, bold=True, color=WHITE, font="Calibri")
txt(sl, "Complete Technical Deep-Dive", Inches(1.0), Inches(3.72), Inches(10), Inches(0.5),
    size=20, color=GREEN, font="Calibri")
txt(sl, "Options Pricing  ·  Monte Carlo Simulation  ·  GBM  ·  Black-Scholes Greeks\nReact Frontend  ·  Three.js 3D  ·  Chart.js Analytics  ·  Vercel Deployment",
    Inches(1.0), Inches(4.25), Inches(11), Inches(0.85),
    size=13, color=MUTED, font="Calibri")

badge(sl, "QUANTITATIVE FINANCE + FULL STACK", Inches(1.0), Inches(5.3), GREEN)
txt(sl, "github.com/coolaeidek-creator/montecarlo_engine  ·  montecarloengine.vercel.app",
    Inches(1.0), Inches(6.8), Inches(11), Inches(0.4),
    size=10, color=MUTED, font="Courier New")

# ════════════════════════════════════════════════════════════
# SLIDE 2 — TABLE OF CONTENTS
# ════════════════════════════════════════════════════════════
sl = add_slide()
section_tag(sl, "OVERVIEW", BLUE)
slide_title(sl, "What This Deck Covers", "From raw math to live website — every layer explained", BLUE)

cols = [
    ("01  FINANCE BASICS",      ["What is an Option?", "Call vs Put", "Strike Price", "Time to Maturity", "Risk-free Rate"], GREEN),
    ("02  THE MATH ENGINE",     ["Monte Carlo Method", "Geometric Brownian Motion", "GBM Formula & Derivation", "Law of Large Numbers", "Discounting to Present Value"], BLUE),
    ("03  BLACK-SCHOLES",       ["BS Model Overview", "d1 & d2 derivation", "The 5 Greeks explained plainly", "When MC beats BS", "Put-Call Parity"], PURPLE),
    ("04  BACKEND / PYTHON",    ["Project Structure", "Pydantic Models", "OptionPricer Class", "NumPy Vectorization", "Confidence Intervals"], ORANGE),
    ("05  FRONTEND / REACT",    ["React 18 State & Hooks", "CDN-only Architecture", "GSAP Animations", "Three.js 3D Scene", "Chart.js Analytics"], PINK),
    ("06  DEPLOYMENT",          ["Git & GitHub", "Vercel Static Deploy", "vercel.json Config", "CI/CD Auto-deploy", "Live URL"], GREEN),
]

col_w = Inches(2.05)
for idx, (title, items, color) in enumerate(cols):
    col = idx % 3
    row = idx // 3
    cx = Inches(0.35) + col * Inches(4.35)
    cy = Inches(2.1) + row * Inches(2.3)
    info_card(sl, title, items, cx, cy, col_w * 2.05, Inches(2.1), color)

# ════════════════════════════════════════════════════════════
# SLIDE 3 — WHAT IS AN OPTION?
# ════════════════════════════════════════════════════════════
sl = add_slide()
section_tag(sl, "FINANCE BASICS", GREEN)
slide_title(sl, "What Is a Stock Option?", "A right — not obligation — to buy or sell a stock at a fixed price")

# analogy box
shape = sl.shapes.add_shape(5, Inches(0.45), Inches(2.1), Inches(5.8), Inches(1.65))
shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x00, 0x20, 0x10)
shape.line.color.rgb = GREEN; shape.line.width = Pt(1); shape.adjustments[0] = 0.05
txt(sl, "Real-Life Analogy", Inches(0.65), Inches(2.18), Inches(5), Inches(0.3),
    size=11, bold=True, color=GREEN)
txt(sl, 'Imagine a house is listed for $500,000. You pay $5,000 to "lock in" that price for\n'
        '3 months. If the house jumps to $600,000, you exercise your right and save $95,000.\n'
        'If it drops, you walk away — losing only your $5,000 deposit.',
    Inches(0.65), Inches(2.5), Inches(5.5), Inches(1.1), size=11, color=MUTED)

# two cards
for i, (label, title, body, color) in enumerate([
    ("CALL OPTION", "Right to BUY",
     ["You profit when stock goes UP", "Pay a premium upfront", "Max loss = premium paid", "Unlimited upside potential", "Used when you're BULLISH"], GREEN),
    ("PUT OPTION", "Right to SELL",
     ["You profit when stock goes DOWN", "Pay a premium upfront", "Max loss = premium paid", "Downside is protected", "Used when you're BEARISH or hedging"], ORANGE),
]):
    cx = Inches(6.6) + i * Inches(3.3)
    info_card(sl, f"{label}  —  {title}", body, cx, Inches(2.1), Inches(3.1), Inches(2.6), color)

txt(sl, "Key Terms at a Glance", Inches(0.45), Inches(3.95), Inches(5), Inches(0.3),
    size=12, bold=True, color=WHITE)
terms = [
    ("Strike Price (K)",     "The agreed price at which you can buy/sell the stock"),
    ("Premium",              "The cost of buying the option contract"),
    ("Expiry / Maturity (T)","The deadline by which you must exercise the option"),
    ("In-The-Money (ITM)",   "The option has intrinsic value right now"),
    ("Out-of-The-Money (OTM)","The option has no intrinsic value — only hope!"),
]
cy = Inches(4.3)
for term, defn in terms:
    txt(sl, f"• {term}:", Inches(0.55), cy, Inches(2.4), Inches(0.3), size=10, bold=True, color=WHITE)
    txt(sl, defn, Inches(2.85), cy, Inches(9.0), Inches(0.3), size=10, color=MUTED)
    cy += Inches(0.32)

# ════════════════════════════════════════════════════════════
# SLIDE 4 — MONTE CARLO METHOD
# ════════════════════════════════════════════════════════════
sl = add_slide()
section_tag(sl, "THE MATH ENGINE", BLUE)
slide_title(sl, "Monte Carlo Simulation", "Pricing options by simulating thousands of random futures", BLUE)

# What is it
shape = sl.shapes.add_shape(5, Inches(0.45), Inches(2.1), Inches(5.9), Inches(1.5))
shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x00, 0x10, 0x25)
shape.line.color.rgb = BLUE; shape.line.width = Pt(1); shape.adjustments[0] = 0.05
txt(sl, "The Core Idea", Inches(0.65), Inches(2.18), Inches(5), Inches(0.3),
    size=11, bold=True, color=BLUE)
txt(sl, "Roll a dice 50,000 times. Record every outcome. Average them.\n"
        "That average, discounted back to today, is your option price.\n"
        "More rolls = more accurate. This is Monte Carlo.",
    Inches(0.65), Inches(2.5), Inches(5.6), Inches(1.0), size=11, color=MUTED)

# steps
steps = [
    ("1", "Generate random shock",   "Draw Z ~ N(0,1) using Box-Muller transform",           GREEN),
    ("2", "Simulate future price",   "Apply GBM formula to get terminal stock price S_T",    BLUE),
    ("3", "Calculate payoff",        "Call: max(S_T - K, 0)   Put: max(K - S_T, 0)",        PURPLE),
    ("4", "Discount to today",       "Multiply payoff by e^(-rT) — the time-value factor",   ORANGE),
    ("5", "Average all payoffs",     "E[Payoff] over N paths = your Monte Carlo price",      PINK),
]
cy = Inches(2.1)
cx = Inches(6.6)
for num, step, detail, color in steps:
    circ = sl.shapes.add_shape(9, cx, cy + Inches(0.04), Inches(0.34), Inches(0.34))
    circ.fill.solid(); circ.fill.fore_color.rgb = color; circ.line.fill.background()
    txt(sl, num, cx, cy + Inches(0.04), Inches(0.34), Inches(0.34),
        size=11, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(sl, step,   cx + Inches(0.42), cy,               Inches(6.1), Inches(0.28), size=11, bold=True,  color=WHITE)
    txt(sl, detail, cx + Inches(0.42), cy + Inches(0.27), Inches(6.1), Inches(0.25), size=10, color=MUTED)
    cy += Inches(0.62)

formula_box(sl, "Price  =  e^(-rT)  ×  (1/N)  ×  Σ  payoff_i",
            Inches(0.45), Inches(3.72), Inches(5.9), Inches(0.55), CARD_BG, BLUE)

txt(sl, "Why Not Just Use a Formula?", Inches(0.45), Inches(4.38), Inches(5.9), Inches(0.3),
    size=11, bold=True, color=WHITE)
bullet(sl, [
    "Closed-form formulas (Black-Scholes) assume too many perfect conditions",
    "Monte Carlo handles path-dependent & exotic options easily",
    "More flexible — just change the payoff function for any option type",
    "Uncertainty is measurable: standard error tells you how accurate you are",
], Inches(0.45), Inches(4.72), Inches(5.9), MUTED, 10)

# ════════════════════════════════════════════════════════════
# SLIDE 5 — GEOMETRIC BROWNIAN MOTION
# ════════════════════════════════════════════════════════════
sl = add_slide()
section_tag(sl, "THE MATH ENGINE", BLUE)
slide_title(sl, "Geometric Brownian Motion (GBM)", "How stocks move — modelled mathematically", BLUE)

formula_box(sl, "S_T  =  S₀ × exp( (r - ½σ²)·T  +  σ·√T·Z )",
            Inches(0.45), Inches(2.05), Inches(12.4), Inches(0.6), CARD_BG, GREEN)

params = [
    ("S₀",     "Current stock price (Spot Price)",         "e.g. $100",   GREEN),
    ("r",      "Risk-free interest rate",                  "e.g. 5%",     BLUE),
    ("σ",      "Volatility — how wild the stock moves",    "e.g. 20%",    PURPLE),
    ("T",      "Time to maturity (in years)",              "e.g. 1.0 yr", ORANGE),
    ("Z",      "Random shock drawn from Normal dist N(0,1)","New per sim", PINK),
    ("½σ²",    "Itô correction — prevents upward bias",   "Math trick",  MUTED),
    ("exp(·)", "Ensures price can never go negative",      "Stock can't = -$5", WHITE),
]
cy = Inches(2.85)
for sym, meaning, example, color in params:
    sym_box = sl.shapes.add_shape(5, Inches(0.45), cy, Inches(0.7), Inches(0.32))
    sym_box.fill.solid(); sym_box.fill.fore_color.rgb = CARD2_BG
    sym_box.line.color.rgb = color; sym_box.line.width = Pt(0.8); sym_box.adjustments[0] = 0.1
    txt(sl, sym, Inches(0.45), cy, Inches(0.7), Inches(0.32),
        size=11, bold=True, color=color, align=PP_ALIGN.CENTER, font="Courier New")
    txt(sl, meaning,  Inches(1.25), cy,              Inches(8.0), Inches(0.3), size=11, color=WHITE)
    txt(sl, example,  Inches(9.4),  cy,              Inches(3.0), Inches(0.3), size=10, color=MUTED, align=PP_ALIGN.RIGHT)
    cy += Inches(0.38)

# drift vs diffusion
txt(sl, "Two Forces That Drive Stock Prices", Inches(0.45), Inches(5.6), Inches(12), Inches(0.3),
    size=12, bold=True, color=WHITE)
info_card(sl, "DRIFT  →  (r - ½σ²)·T",
          ["Deterministic upward trend", "Based on risk-free rate r", "The ½σ² removes Itô bias", "This is the 'expected direction'"],
          Inches(0.45), Inches(5.95), Inches(6.0), Inches(1.35), BLUE)
info_card(sl, "DIFFUSION  →  σ·√T·Z",
          ["Random component — unpredictability", "Higher σ = wilder price swings", "√T: uncertainty grows with time", "Z is freshly random each simulation"],
          Inches(6.6), Inches(5.95), Inches(6.1), Inches(1.35), PURPLE)

# ════════════════════════════════════════════════════════════
# SLIDE 6 — BOX-MULLER & STATISTICS
# ════════════════════════════════════════════════════════════
sl = add_slide()
section_tag(sl, "THE MATH ENGINE", BLUE)
slide_title(sl, "Generating Randomness & Measuring Accuracy", "Box-Muller transform · Standard Error · Confidence Intervals", BLUE)

# Box-Muller
txt(sl, "Box-Muller Transform — Making Normal Random Numbers", Inches(0.45), Inches(2.1), Inches(8), Inches(0.3),
    size=13, bold=True, color=GREEN)
txt(sl, "Computers generate uniform random numbers (0 to 1). We need bell-curve (Normal) numbers.\nBox-Muller converts two uniform numbers (U₁, U₂) into one Normal number Z:",
    Inches(0.45), Inches(2.45), Inches(8), Inches(0.5), size=11, color=MUTED)
formula_box(sl, "Z  =  √(-2·ln U₁)  ×  cos(2π·U₂)",
            Inches(0.45), Inches(3.0), Inches(5.5), Inches(0.52), CARD_BG, GREEN)
code_block(sl,
    "function randn() {\n"
    "  let u = 0, v = 0;\n"
    "  while (!u) u = Math.random();\n"
    "  while (!v) v = Math.random();\n"
    "  return Math.sqrt(-2*Math.log(u)) * Math.cos(2*Math.PI*v);\n"
    "}",
    Inches(0.45), Inches(3.62), Inches(5.5), Inches(1.2))

# Right: SE & CI
txt(sl, "How Accurate Is Our Price? — Standard Error", Inches(6.6), Inches(2.1), Inches(6.3), Inches(0.3),
    size=13, bold=True, color=ORANGE)
bullet(sl, [
    "Run N simulations → get N payoff values",
    "Average them → that's your price estimate",
    "Standard deviation of payoffs ÷ √N = Standard Error",
    "More simulations = smaller error = more confidence",
], Inches(6.6), Inches(2.48), Inches(6.1), MUTED, 11)
formula_box(sl, "SE  =  std(payoffs) / √N",
            Inches(6.6), Inches(3.72), Inches(6.1), Inches(0.5), CARD_BG, ORANGE)

txt(sl, "95% Confidence Interval", Inches(6.6), Inches(4.35), Inches(6.1), Inches(0.28),
    size=12, bold=True, color=WHITE)
txt(sl, "We are 95% confident the true price lies between:",
    Inches(6.6), Inches(4.65), Inches(6.1), Inches(0.28), size=11, color=MUTED)
formula_box(sl, "CI  =  [ Price - 1.96×SE ,  Price + 1.96×SE ]",
            Inches(6.6), Inches(5.0), Inches(6.1), Inches(0.5), CARD_BG, BLUE)
txt(sl, "1.96 comes from the Normal distribution — covers 95% of outcomes.\n"
        "With 50,000 simulations our SE is typically < $0.03.",
    Inches(6.6), Inches(5.6), Inches(6.1), Inches(0.5), size=10, color=MUTED)

txt(sl, "Law of Large Numbers in Action", Inches(0.45), Inches(4.92), Inches(5.5), Inches(0.28),
    size=12, bold=True, color=WHITE)
bullet(sl, [
    "1,000 paths  → rough estimate, SE ≈ $0.25",
    "10,000 paths → good estimate, SE ≈ $0.08",
    "50,000 paths → precise,      SE ≈ $0.03",
    "100,000 paths→ very precise, SE ≈ $0.02",
    "As N → ∞, MC price converges to true price",
], Inches(0.45), Inches(5.25), Inches(5.5), MUTED, 10)

# ════════════════════════════════════════════════════════════
# SLIDE 7 — BLACK-SCHOLES & GREEKS
# ════════════════════════════════════════════════════════════
sl = add_slide()
section_tag(sl, "BLACK-SCHOLES MODEL", PURPLE)
slide_title(sl, "Black-Scholes & The Option Greeks", "Analytical formulas vs Monte Carlo — and what Δ Γ Θ ν really mean", PURPLE)

formula_box(sl, "d₁ = [ln(S/K) + (r + ½σ²)T] / (σ√T)         d₂ = d₁ - σ√T",
            Inches(0.45), Inches(2.05), Inches(12.4), Inches(0.52), CARD_BG, PURPLE)

greeks = [
    ("Δ  DELTA", "If stock goes up $1, how much does my option gain?",
     "Call Delta = 0.64 → option gains $0.64 per $1 stock move\nPut Delta = -0.36 → option loses $0.36 per $1 stock move\nAtm option has delta ≈ 0.5  |  Deep ITM ≈ 1.0",
     GREEN),
    ("Γ  GAMMA", "How fast is Delta itself changing?",
     "High gamma = option is very sensitive (near expiry, near strike)\nLike acceleration vs speed — delta is speed, gamma is acceleration\nATM options have highest gamma",
     BLUE),
    ("Θ  THETA", "How much value does the option lose every single day?",
     'Options are "melting ice cubes" — time kills their value\nTheta = -$0.017 means you lose $0.017 of value overnight\nSellers love theta; buyers must be right fast',
     ORANGE),
    ("ν  VEGA", "If volatility spikes 1%, how much does my option gain?",
     "Vega = $0.38 → option gains $0.38 if vol goes from 20% to 21%\nHigh vega = option is very sensitive to market panic or calm\nBuy options before earnings (vol expected to spike)",
     PURPLE),
]
cy = Inches(2.72)
for sym, q, ans, color in greeks:
    info_card(sl, f"{sym}  —  {q}", ans.split('\n'), Inches(0.45), cy, Inches(12.4), Inches(1.08), color)
    cy += Inches(1.15)

# ════════════════════════════════════════════════════════════
# SLIDE 8 — PUT-CALL PARITY
# ════════════════════════════════════════════════════════════
sl = add_slide()
section_tag(sl, "FINANCIAL MATH", GREEN)
slide_title(sl, "Put-Call Parity — The Sanity Check", "A fundamental no-arbitrage relationship between calls and puts")

formula_box(sl, "C - P  =  S₀ - K · e^(-rT)",
            Inches(2.0), Inches(2.05), Inches(9.0), Inches(0.62), CARD_BG, GREEN)

txt(sl, "In Plain English:", Inches(0.45), Inches(2.85), Inches(12), Inches(0.3),
    size=13, bold=True, color=WHITE)
txt(sl, '(Call Price) minus (Put Price) must equal (Spot Price) minus (Present Value of Strike).\n'
        'If this breaks, traders can make risk-free profit — so the market forces it to hold.',
    Inches(0.45), Inches(3.2), Inches(12), Inches(0.5), size=12, color=MUTED)

sides = [
    ("LEFT SIDE:  C - P",
     ["C = Call option price (right to buy)", "P = Put option price (right to sell)",
      "The difference between the two", "Should equal the right side exactly"],
     GREEN),
    ("RIGHT SIDE:  S₀ - K·e^(-rT)",
     ["S₀ = Current stock price today", "K = Strike price of both options",
      "e^(-rT) = discount factor (time value of money)", "Present value of owning the stock vs paying strike later"],
     BLUE),
]
for i, (title, items, color) in enumerate(sides):
    info_card(sl, title, items, Inches(0.45) + i*Inches(6.5), Inches(3.82), Inches(6.1), Inches(1.75), color)

txt(sl, "Why We Check It:", Inches(0.45), Inches(5.72), Inches(12), Inches(0.28),
    size=12, bold=True, color=WHITE)
bullet(sl, [
    "Our Monte Carlo uses independent simulations for call and put — parity validates both are correct",
    "Difference < $0.15 = simulation is accurate   |   Difference > $0.30 = need more simulations",
    "This is a standard test used in professional quantitative finance",
], Inches(0.45), Inches(6.05), Inches(12), MUTED, 11)

# ════════════════════════════════════════════════════════════
# SLIDE 9 — PYTHON BACKEND
# ════════════════════════════════════════════════════════════
sl = add_slide()
section_tag(sl, "BACKEND — PYTHON", ORANGE)
slide_title(sl, "Python Backend Architecture", "Clean separation of concerns: models → simulator → pricer", ORANGE)

# Project tree
code_block(sl,
    "montecarlo_engine/\n"
    "├── engine/\n"
    "│   ├── models.py      # Data structures (Pydantic)\n"
    "│   ├── random.py      # N(0,1) random generation\n"
    "│   ├── simulator.py   # GBM terminal price simulation\n"
    "│   ├── payoff.py      # Call / Put payoff calculation\n"
    "│   └── pricer.py      # OptionPricer class (orchestrator)\n"
    "├── app/\n"
    "│   └── main.py        # Entry point — runs & prints results\n"
    "├── index.html         # Full frontend (React + Three.js)\n"
    "└── vercel.json        # Deployment config",
    Inches(0.45), Inches(2.08), Inches(5.5), Inches(2.8))

# Modules explained
modules = [
    ("models.py — Pydantic Data Classes",
     ["MarketEnvironment: spot, rate, volatility, maturity",
      "OptionContract: strike price, option type (call/put)",
      "Pydantic validates inputs automatically — no bad data",
      "Type hints make code self-documenting"],
     GREEN),
    ("simulator.py — GBM in NumPy",
     ["simulate_terminal_prices(market, shocks)",
      "Takes N(0,1) shocks, applies GBM formula",
      "NumPy vectorized — runs 50k sims in milliseconds",
      "Returns array of N terminal prices S_T"],
     BLUE),
    ("pricer.py — The Orchestrator",
     ["OptionPricer class with price() method",
      "Calls random → simulator → payoff in sequence",
      "Discounts payoffs by e^(-rT)",
      "Returns price, std_error, confidence_interval dict"],
     PURPLE),
]
cy = Inches(2.08)
for title, items, color in modules:
    info_card(sl, title, items, Inches(6.2), cy, Inches(6.7), Inches(1.42), color)
    cy += Inches(1.52)

txt(sl, "Key Python Concepts Used", Inches(0.45), Inches(5.0), Inches(5.5), Inches(0.28),
    size=11, bold=True, color=WHITE)
bullet(sl, [
    "Pydantic — runtime data validation & type safety",
    "NumPy — vectorized math (50k ops in one line)",
    "np.mean, np.std, np.sqrt — statistical computation",
    "Dataclasses / OOP — clean, testable, reusable code",
], Inches(0.45), Inches(5.32), Inches(5.4), MUTED, 10)

# ════════════════════════════════════════════════════════════
# SLIDE 10 — REACT FRONTEND
# ════════════════════════════════════════════════════════════
sl = add_slide()
section_tag(sl, "FRONTEND — REACT", PINK)
slide_title(sl, "React 18 Frontend Architecture", "Zero build step — CDN-only React with Babel standalone", PINK)

# Architecture
arch = [
    ("NO BUILD STEP", "React + Babel + Three.js loaded from CDN URLs in <script> tags. No npm install, no webpack, no node_modules. Browser compiles JSX on the fly.", ORANGE),
    ("STATE MANAGEMENT", "useState hook manages: params (spot/strike/vol/rate/T/N), res (call/put prices), busy (loading state). setParams triggers re-render, runSim recomputes everything.", GREEN),
    ("SIDE EFFECTS", "useEffect runs Three.js scene setup once on mount. Returns cleanup function that destroys WebGL context when component unmounts. Proper memory management.", BLUE),
    ("PERFORMANCE", "useCallback memoizes runSim so it doesn't recreate on every render. setTimeout(fn, 30) prevents UI freeze during Monte Carlo loop. Chart.js destroyed & recreated on each sim.", PURPLE),
]
cy = Inches(2.08)
for i, (title, body, color) in enumerate(arch):
    col = i % 2
    row = i // 2
    cx = Inches(0.45) + col * Inches(6.5)
    cy2 = Inches(2.08) + row * Inches(1.52)
    info_card(sl, title, [body[:80], body[80:] if len(body) > 80 else ""], cx, cy2, Inches(6.1), Inches(1.38), color)

txt(sl, "Component Tree", Inches(0.45), Inches(5.22), Inches(12.4), Inches(0.28),
    size=12, bold=True, color=WHITE)
code_block(sl,
    "<App>  — manages all state, runs simulation, owns scene3d ref\n"
    "  <Slider>       — reusable range input with fill bar & live value\n"
    "  <AnimNum>      — GSAP-animated number that tweens from old → new value\n"
    "  <ChartPanel>   — 3 Chart.js charts in tabbed view (Payoff/Dist/Paths)\n"
    "  <SuggestionsPanel>  — 5 analysis cards using computed Greeks & signals",
    Inches(0.45), Inches(5.55), Inches(12.4), Inches(1.62))

# ════════════════════════════════════════════════════════════
# SLIDE 11 — THREE.JS 3D SCENE
# ════════════════════════════════════════════════════════════
sl = add_slide()
section_tag(sl, "FRONTEND — THREE.JS", BLUE)
slide_title(sl, "Three.js 3D Background Scene", "WebGL-powered live 3D canvas — runs at 60fps behind the UI", BLUE)

objects = [
    ("⭐ STAR FIELD",     "4,000 random points in 3D space. BufferGeometry with Float32Array positions. PointsMaterial for tiny white dots. Slowly rotates on X and Y axes.", GREEN),
    ("◈ WIREFRAME ORBS",  "3 IcosahedronGeometry meshes (10r, 6r, 3r). MeshBasicMaterial wireframe=true. Each rotates at different speeds creating depth illusion.", BLUE),
    ("○ TORUS RING",      "TorusGeometry(14, 0.06) — thin glowing ring. Tilted at π/2.8 radians. Rotates on Z axis. Represents the 'orbital' theme.", PURPLE),
    ("● COLORED ORBS",    "7 small spheres in brand colors. Float up/down using sin() wave per orb. Each has a random phase offset so they don't sync.", ORANGE),
    ("〜 GBM PATH TUBES",  "22 CatmullRomCurve3 paths through 3D space. TubeGeometry wraps each curve. Rebuilt on every simulation with new random GBM data. Colors match brand palette.", PINK),
    ("📷 CAMERA DRIFT",    "Camera orbits lazily: position.x = sin(t×0.065)×4.5. position.y = sin(t×0.048)×2.8. lookAt(0,0,0) always. Creates cinematic depth with no user input.", WHITE),
]
cy = Inches(2.08)
for i, (title, body, color) in enumerate(objects):
    col = i % 2
    row = i // 2
    cx = Inches(0.45) + col * Inches(6.5)
    cy2 = Inches(2.08) + row * Inches(1.52)
    info_card(sl, title, [body[:85], body[85:] if len(body) > 85 else ""], cx, cy2, Inches(6.1), Inches(1.38), color)

txt(sl, "Fog: FogExp2(0x030712, 0.016) — distant objects fade into background · FPS: requestAnimationFrame loop · Pixel ratio: min(devicePixelRatio, 2)",
    Inches(0.45), Inches(6.85), Inches(12.4), Inches(0.35), size=9, color=MUTED)

# ════════════════════════════════════════════════════════════
# SLIDE 12 — GSAP ANIMATIONS
# ════════════════════════════════════════════════════════════
sl = add_slide()
section_tag(sl, "FRONTEND — GSAP + CHART.JS", PINK)
slide_title(sl, "GSAP Animations & Chart.js Analytics", "Smooth motion design + interactive data visualisation", PINK)

txt(sl, "GSAP — GreenSock Animation Platform", Inches(0.45), Inches(2.08), Inches(6.1), Inches(0.3),
    size=13, bold=True, color=PINK)
gsap_items = [
    ("Entrance animations", "Header slides from top (y:-44→0). Control panel slides from left. Results/charts slide from right. Staggered delays = professional feel."),
    ("AnimNum component",    "gsap.to(obj, {v: newValue}) tweens a plain JS object. onUpdate writes to DOM via ref. Number smoothly counts from old price to new price."),
    ("Card reveal",          "gsap.fromTo('.opt-card', {y:10, opacity:0}, {y:0, opacity:1, stagger:0.07}) — each card enters with a slight cascade delay."),
    ("Button feedback",      "gsap.to('.run-btn', {scale:0.95, yoyo:true, repeat:1}) — micro-interaction on click. Feels physical and responsive."),
]
cy = Inches(2.45)
for title, body in gsap_items:
    txt(sl, f"• {title}:", Inches(0.55), cy, Inches(2.8), Inches(0.3), size=10, bold=True, color=PINK)
    txt(sl, body, Inches(3.25), cy, Inches(3.2), Inches(0.42), size=10, color=MUTED)
    cy += Inches(0.46)

accent_line(sl, Inches(6.55), Inches(2.08), Inches(0.04), BLUE, h=Inches(4.8))

txt(sl, "Chart.js — 3 Interactive Chart Types", Inches(6.7), Inches(2.08), Inches(6.2), Inches(0.3),
    size=13, bold=True, color=BLUE)
charts_info = [
    ("PAYOFF CHART (Line)",
     ["Call & put P&L curves across all possible stock prices",
      "X-axis: terminal price  Y-axis: P&L in dollars",
      "Custom vlinePlugin draws vertical strike & spot markers",
      "Zero-line shows break-even visually"], BLUE),
    ("DISTRIBUTION (Bar)",
     ["6,000 simulated terminal prices binned into histogram",
      "Green bars: ITM for call  |  Orange bars: ITM for put",
      "Shows the probability distribution of outcomes",
      "You can literally see your odds of profit"], GREEN),
    ("GBM PATHS (Line)",
     ["22 individual simulated stock price journeys",
      "Each path is a full GBM simulation with 80 time steps",
      "Shows the 'fan' of possible futures",
      "Strike shown as dashed line — you can see ITM vs OTM"], PURPLE),
]
cy2 = Inches(2.48)
for title, items, color in charts_info:
    info_card(sl, title, items, Inches(6.7), cy2, Inches(6.25), Inches(1.32), color)
    cy2 += Inches(1.42)

# ════════════════════════════════════════════════════════════
# SLIDE 13 — DEPLOYMENT
# ════════════════════════════════════════════════════════════
sl = add_slide()
section_tag(sl, "DEPLOYMENT", GREEN)
slide_title(sl, "Git, GitHub & Vercel Deployment", "From local code to live website in under 30 seconds", GREEN)

flow = ["Write Code", "git add", "git commit", "git push", "vercel deploy", "Live on Web!"]
colors_f = [WHITE, GREEN, GREEN, BLUE, PURPLE, ORANGE]
cx = Inches(0.45)
for i, (step, color) in enumerate(zip(flow, colors_f)):
    box_s = sl.shapes.add_shape(5, cx, Inches(2.08), Inches(2.0), Inches(0.52))
    box_s.fill.solid(); box_s.fill.fore_color.rgb = CARD_BG
    box_s.line.color.rgb = color; box_s.line.width = Pt(1); box_s.adjustments[0] = 0.25
    txt(sl, step, cx, Inches(2.08), Inches(2.0), Inches(0.52),
        size=10, bold=True, color=color, align=PP_ALIGN.CENTER)
    cx += Inches(2.15)
    if i < len(flow)-1:
        txt(sl, "→", cx - Inches(0.22), Inches(2.12), Inches(0.3), Inches(0.38),
            size=14, color=MUTED, align=PP_ALIGN.CENTER)

cards_dep = [
    ("GIT — VERSION CONTROL",
     ["git init — start tracking your project",
      "git add <file> — stage specific changes",
      "git commit -m 'message' — save a checkpoint",
      "git push — upload to GitHub cloud",
      "git log — see full history of changes",
      "Every commit has: author, timestamp, message"],
     GREEN),
    ("VERCEL — STATIC HOSTING",
     ["vercel.json: sets framework=null, outputDirectory='.'",
      "Tells Vercel: 'serve index.html as static site'",
      "vercel deploy --prod: uploads files to CDN",
      "Aliased to montecarloengine.vercel.app",
      "Edge network: fast worldwide delivery",
      "Free tier: unlimited static deploys"],
     BLUE),
    ("GITHUB — CODE HOSTING",
     ["Remote repository: stores code in the cloud",
      "gh auth login: authenticate CLI to GitHub",
      "gh auth setup-git: enables HTTPS push via token",
      "git config user.email: ties commits to your account",
      "Public repo: anyone can view & fork the code",
      "Connected to Vercel for auto-deploy on push"],
     PURPLE),
    ("ARCHITECTURE CHOICES",
     ["No Node.js needed: CDN-only frontend",
      "Single file: entire app = one index.html",
      "No database: all computation in-browser JS",
      "No server: pure static = zero hosting cost",
      "No build pipeline: ship HTML directly",
      "Babel standalone: JSX compiled in browser"],
     ORANGE),
]
cy = Inches(2.75)
for i, (title, items, color) in enumerate(cards_dep):
    col = i % 2
    row = i // 2
    cx2 = Inches(0.45) + col * Inches(6.5)
    cy2 = Inches(2.75) + row * Inches(2.15)
    info_card(sl, title, items, cx2, cy2, Inches(6.1), Inches(2.0), color)

# ════════════════════════════════════════════════════════════
# SLIDE 14 — SUGGESTIONS ENGINE
# ════════════════════════════════════════════════════════════
sl = add_slide()
section_tag(sl, "ANALYSIS ENGINE", PURPLE)
slide_title(sl, "The Suggestions & Analysis Engine", "5 smart cards computed from your simulation results", PURPLE)

suggestions = [
    ("MONEYNESS GAUGE",
     "Compares Spot vs Strike: (S/K - 1) × 100\n"
     "< -15% → Deep OTM (needs big move)  |  ±2% → ATM (highest gamma)\n"
     "> +15% → Deep ITM (behaves like stock). Needle on gradient bar shows position.",
     GREEN),
    ("OPTION SENSITIVITY (GREEKS)",
     "Uses Black-Scholes closed-form to compute Δ,Γ,Θ,ν in real time.\n"
     "Displayed in plain English: 'If stock goes up $1, call gains $X'.\n"
     "normalCDF and normalPDF implemented from scratch (Abramowitz & Stegun approx).",
     BLUE),
    ("BREAKEVEN PRICES",
     "Call breakeven = Strike + Call Premium\n"
     "Put  breakeven = Strike - Put  Premium\n"
     "Shows % move needed from current spot. Trader's most important number.",
     ORANGE),
    ("STRATEGY SIGNAL",
     "Rule-based engine using vol regime + moneyness:\n"
     "Vol > 50% → Sell Straddle  |  Vol < 15% → Buy Straddle\n"
     "ITM > 10% → Covered Call  |  OTM > 10% → Bull Call Spread\n"
     "Near ATM → Iron Condor. Each has plain-English reasoning.",
     PINK),
    ("RISK SCORE 1–10",
     "Composite score: vol × 11 + |moneyness|/14 + maturity × 0.4\n"
     "Clamped to 1–10 range. Color: green=low, blue=medium, orange=high, pink=extreme.\n"
     "10 bar segments fill up with risk color. Instantly readable risk level.",
     PURPLE),
]
cy = Inches(2.08)
for i, (title, body, color) in enumerate(suggestions):
    lines = body.split('\n')
    info_card(sl, title, lines, Inches(0.45), cy, Inches(12.4), Inches(0.98), color)
    cy += Inches(1.06)

# ════════════════════════════════════════════════════════════
# SLIDE 15 — FINAL SUMMARY
# ════════════════════════════════════════════════════════════
sl = add_slide()
section_tag(sl, "SUMMARY", GREEN)
slide_title(sl, "Everything At A Glance", "The full stack — finance → math → code → deployment")

cols_s = [
    ("FINANCE LAYER",
     ["European Call & Put options", "Intrinsic value vs time value",
      "Moneyness: ITM/ATM/OTM", "Put-Call Parity no-arbitrage",
      "Strategy signals & breakevens"],
     GREEN),
    ("MATH LAYER",
     ["Monte Carlo simulation (N=50k)", "Geometric Brownian Motion (GBM)",
      "Box-Muller Normal sampling", "Itô Lemma drift correction",
      "Black-Scholes Greeks: Δ,Γ,Θ,ν"],
     BLUE),
    ("PYTHON BACKEND",
     ["Pydantic models for validation", "NumPy vectorized simulation",
      "OptionPricer OOP class", "Standard error & 95% CI",
      "PYTHONPATH module structure"],
     PURPLE),
    ("REACT FRONTEND",
     ["useState / useEffect / useRef hooks", "GSAP number tweening & entrances",
      "Three.js WebGL 3D scene at 60fps", "Chart.js: Payoff, Dist, Paths charts",
      "CDN-only, zero-build architecture"],
     ORANGE),
    ("DEPLOYMENT STACK",
     ["Git version control", "GitHub remote repository",
      "gh CLI authentication & push", "vercel.json static site config",
      "montecarloengine.vercel.app live"],
     PINK),
]
col_w = Inches(2.5)
for i, (title, items, color) in enumerate(cols_s):
    cx = Inches(0.35) + i * Inches(2.6)
    info_card(sl, title, items, cx, Inches(2.08), col_w, Inches(4.4), color)

accent_line(sl, Inches(0.45), Inches(6.65), Inches(12.4), GREEN)
txt(sl, "github.com/coolaeidek-creator/montecarlo_engine  ·  Live: montecarloengine.vercel.app",
    Inches(0.45), Inches(6.72), Inches(12.4), Inches(0.4),
    size=11, color=MUTED, align=PP_ALIGN.CENTER, font="Courier New")

# ── Save ──────────────────────────────────────────────────
out = "/Users/ronan/CoolAeidek/montecarlo_engine/MonteCarlo_Engine_Deep_Dive.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
